from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# Define colors
DARK_BG = RGBColor(15, 23, 42)      # #0f172a
LIGHT_BG = RGBColor(30, 41, 59)     # #1e293b
BLUE = RGBColor(59, 130, 246)       # #3b82f6
GREEN = RGBColor(34, 197, 94)       # #22c55e
YELLOW = RGBColor(251, 191, 36)     # #fbbf24
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(148, 163, 184)      # #94a3b8
DARK_GRAY = RGBColor(100, 116, 139) # #64748b
RED = RGBColor(220, 38, 38)         # #dc2626

def add_title_slide(prs, title, subtitle, tagline, stats):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.5))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Stats
    stat_width = 3
    start_x = (13.333 - len(stats) * stat_width) / 2
    for i, (num, label) in enumerate(stats):
        x = start_x + i * stat_width
        # Number
        num_box = slide.shapes.add_textbox(Inches(x), Inches(5), Inches(stat_width), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(5.8), Inches(stat_width), Inches(0.5))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_left, content_right=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    
    # Header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = LIGHT_BG
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Blue accent line
    line = slide.shapes.add_shape(1, Inches(0), Inches(0.8), Inches(13.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    
    if content_right:
        # Two column layout
        col_width = 6
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(col_width), Inches(6))
        tf = left_box.text_frame
        tf.word_wrap = True
        for item in content_left:
            if isinstance(item, tuple):
                title_text, desc = item
                p = tf.add_paragraph()
                p.text = title_text
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = YELLOW
                p.space_before = Pt(8)
                p = tf.add_paragraph()
                p.text = desc
                p.font.size = Pt(11)
                p.font.color.rgb = GRAY
            else:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(7), Inches(1), Inches(col_width), Inches(6))
        tf = right_box.text_frame
        tf.word_wrap = True
        for item in content_right:
            if isinstance(item, tuple):
                title_text, desc = item
                p = tf.add_paragraph()
                p.text = title_text
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = GREEN
                p.space_before = Pt(8)
                p = tf.add_paragraph()
                p.text = desc
                p.font.size = Pt(11)
                p.font.color.rgb = GRAY
            else:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.333), Inches(6))
        tf = content_box.text_frame
        tf.word_wrap = True
        for item in content_left:
            if isinstance(item, tuple):
                title_text, desc = item
                p = tf.add_paragraph()
                p.text = title_text
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = YELLOW
                p.space_before = Pt(10)
                p = tf.add_paragraph()
                p.text = desc
                p.font.size = Pt(11)
                p.font.color.rgb = GRAY
            else:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
    
    return slide

# Slide 1: Title
add_title_slide(prs, "PairProxy", "企业级 LLM API 代理网关", 
                "统一管控 · 精确追踪 · 零侵入接入",
                [("57,550+", "代码行数"), ("1,100+", "测试用例"), ("99.5%", "可用性")])

# Slide 2: Pain Points & Solutions
pain_points = [
    ("API Key 泄露风险", "分发给每个人，无法追溯责任"),
    ("用量不透明", "月底看账单才知道超支"),
    ("无法限流", "某人用掉大量资源无法拦截"),
    ("费用分摊困难", "无法精确计算每人消耗")
]
solutions = [
    ("集中存储，用户不可见", "真实 API Key 仅存储在服务端"),
    ("实时统计，Dashboard 可视化", "费用估算精确到小数点后6位"),
    ("配额检查，超额返回 429", "按分组设置日/月 token 上限"),
    ("按用户统计，支持导出", "精确追踪每人用量")
]
add_content_slide(prs, "业务痛点与解决方案", pain_points, solutions)

# Slide 3: Architecture
arch_content = [
    ("系统架构", "Claude Code → cproxy(本地:8080) → sproxy(服务端:9000) → Anthropic/OpenAI API"),
    "",
    ("JWT 认证", "Access 24h + Refresh 7天，自动刷新"),
    ("配额管理", "日/月 token 上限，RPM 滑动窗口限流"),
    ("负载均衡", "两级负载均衡，健康检查"),
    ("集群模式", "Primary + Worker，自动路由"),
    ("Token 统计", "精确统计 input/output，零延迟"),
    ("零侵入接入", "用户仅需设置两个环境变量")
]
add_content_slide(prs, "系统架构：两层代理设计", arch_content)

# Slide 4: Capabilities
caps = [
    ("认证与授权", "JWT Token 自动刷新、LDAP/AD 集成、登录频率限制、Token 黑名单"),
    ("配额与限流", "日/月 Token 上限、RPM 滑动窗口限流、分组配额管理、超额自动拦截"),
    ("负载与高可用", "两级负载均衡、主动+被动健康检查、请求级自动重试、集群模式支持"),
    ("可观测性", "Prometheus 指标、Webhook 告警、OpenTelemetry 追踪、管理审计日志"),
    ("运维管理", "Web Dashboard、Admin CLI 工具、数据导出/备份、对话内容追踪"),
    ("安全合规", "bcrypt 密码加密、强制 HS256 算法、LDAPS 支持、审计日志留存")
]
add_content_slide(prs, "核心能力矩阵", caps)

# Slide 5: Value & ROI
value_content = [
    ("成本透明", "实时追踪每笔请求费用，精确到 0.000001 USD"),
    ("预算可控", "按团队/个人设置日/月配额，超额自动拦截"),
    ("责任可追", "完整审计日志，谁用了多少一目了然"),
    ("安全合规", "API Key 集中管控，支持 LDAP/AD"),
    "",
    "投资回报预估（50人团队年度收益）：",
    "• 30% 用量优化",
    "• 100% 成本可视",
    "• 99.5% 服务可用性",
    "• 0 额外运维成本"
]
add_content_slide(prs, "业务价值与 ROI", value_content)

# Slide 6: AI Efficiency
ai_content = [
    "由 Claude AI 辅助开发，展现企业级软件开发新范式",
    "代码质量达到生产就绪标准，完全可交付使用",
    "",
    ("57,550+", "代码行数 Go 1.24"),
    ("1,100+", "测试用例 全部通过"),
    ("~70%", "测试覆盖率 核心 80%+"),
    ("11+", "技术文档 200KB+"),
    "",
    ("22", "内部模块 高度解耦"),
    ("100%", "测试通过率 无数据竞争"),
    ("99.5%", "系统可用性 Silver 等级"),
    ("v2.5.0", "生产就绪 已验收通过")
]
add_content_slide(prs, "AI Coding 开发效率展示", ai_content)

# Slide 7: Production Readiness
prod_content = [
    ("质量保证体系", ""),
    "• 测试覆盖全面：1,100+ 测试用例，三层覆盖",
    "• 代码质量规范：golangci-lint 通过，~70%覆盖率",
    "• 文档体系完整：11个主要文档，覆盖全场景",
    "",
    ("安全与运维", ""),
    "• 安全模型完备：JWT 防护、bcrypt、LDAPS",
    "• 故障容错分析：24个场景，0个高风险项",
    "",
    ("验收评估结果", ""),
    "• 功能完整性：★★★★★",
    "• 代码质量：★★★★★",
    "• 测试覆盖：★★★★",
    "• 文档完整性：★★★★★",
    "• 生产就绪度：★★★★★",
    "• 综合评分：4.8/5.0"
]
add_content_slide(prs, "产品成熟度：生产就绪评估", prod_content)

# Slide 8: Deployment
deploy_content = [
    ("部署方式", ""),
    "• Docker：多架构镜像，15MB 轻量",
    "• systemd：Linux 生产环境，安全加固",
    "• 二进制：单文件部署，零依赖",
    "• 集群模式：Primary + Worker，自动路由",
    "",
    ("典型应用场景", ""),
    "• 企业 AI 开发团队：50+ 开发者共享，精确追踪用量",
    "• 多 Provider 统一接入：同时管理 Anthropic、OpenAI",
    "• 成本中心管控：IT 统一采购，业务线按配额申请",
    "• 合规审计要求：金融/医疗行业，留存完整对话记录"
]
add_content_slide(prs, "部署方案与应用场景", deploy_content)

# Slide 9: Future & Summary
future_content = [
    ("演进路线图", ""),
    "• P0 - 生产必需（已实现）：水印追踪、请求重试、路由发现",
    "• P1 - 短期优化（1-2月）：OpenClaw 智能运维、自动备份",
    "• P2 - 中期规划（3+月）：Primary HA、分布式追踪",
    "",
    ("项目成果总结", ""),
    "PairProxy 是 AI Coding 开发企业级软件的成功实践：",
    "",
    "1. 代码质量达标：57,550+ 行 Go 代码，1,100+ 测试用例",
    "2. 产品成熟可用：Silver 可用性 99.5%，评分 4.8/5.0",
    "3. 交付标准完整：11+ 文档，CI/CD 自动化，多平台支持",
    "4. AI 开发可行：AI 辅助可完成企业级软件开发"
]
add_content_slide(prs, "未来规划与项目总结", future_content)

# Slide 10: End
add_title_slide(prs, "PairProxy", "企业级 LLM API 代理网关",
                "统一管控 · 精确追踪 · 零侵入接入 · AI 驱动开发",
                [("v2.5.0", "Production Ready"), ("GitHub", "github.com/l17728/pairproxy"), ("License", "Apache 2.0")])

# Save
prs.save('D:\\pairproxy\\PairProxy_Management_Report.pptx')
print("PPT created: D:\\pairproxy\\PairProxy_Management_Report.pptx")
print("Total slides: 10")
